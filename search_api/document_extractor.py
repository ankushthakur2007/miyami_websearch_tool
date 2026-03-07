"""
Document Extractor Module

Extracts text from various document formats and converts to markdown:
- PDF (pypdf)
- Word/DOCX (python-docx)
- Excel/XLSX (openpyxl)
- PowerPoint/PPTX (python-pptx)
- Markdown (raw text)
- RTF (simple text extraction)
- Plain text
"""

import re
from typing import Optional, Dict, Any
from pathlib import Path
import tempfile
import base64


def is_document_url(url: str) -> bool:
    """Check if URL likely points to a document file based on extension."""
    document_extensions = {
        '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
        '.rtf', '.txt', '.odt', '.ods', '.odp', '.csv'
    }
    path = Path(url.lower())
    return path.suffix in document_extensions


def get_content_type_mime(content_type: str) -> Optional[str]:
    """
    Extract document type from Content-Type header.
    Returns document category or None if not a document.
    """
    content_type = content_type.lower().strip()

    mime_to_type = {
        # PDF
        'application/pdf': 'pdf',
        'application/x-pdf': 'pdf',
        'application/acrobat': 'pdf',
        # Word
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'text/rtf': 'rtf',
        # Excel
        'application/vnd.ms-excel': 'xls',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'text/csv': 'csv',
        # PowerPoint
        'application/vnd.ms-powerpoint': 'ppt',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        # Text
        'text/plain': 'txt',
        'text/markdown': 'md',
        'text/richtext': 'rtf',
    }

    return mime_to_type.get(content_type)


def extract_pdf(html_content: bytes) -> str:
    """Extract text from PDF bytes using pypdf."""
    try:
        import pypdf
        import io

        pdf_file = io.BytesIO(html_content)
        reader = pypdf.PdfReader(pdf_file)
        text_parts = []

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

        return "\n\n".join(text_parts) if text_parts else "[No text content found in PDF]"
    except Exception as e:
        return f"[PDF text extraction failed: {str(e)}]"


def extract_docx(html_content: bytes) -> str:
    """Extract text from DOCX bytes using python-docx."""
    try:
        import docx
        import io

        doc_file = io.BytesIO(html_content)
        doc = docx.Document(doc_file)

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        return "\n\n".join(paragraphs) if paragraphs else "[No text content found in document]"
    except Exception as e:
        return f"[DOCX text extraction failed: {str(e)}]"


def extract_excel(html_content: bytes) -> str:
    """Extract text from XLSX bytes using openpyxl."""
    try:
        import openpyxl
        import io

        wb = openpyxl.load_workbook(io.BytesIO(html_content))
        text_parts = []

        for sheet in wb.worksheets:
            sheet_name = f"Sheet: {sheet.title}"
            text_parts.append(sheet_name)
            text_parts.append("-" * len(sheet_name))

            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text_parts.append(row_text)

            text_parts.append("")

        return "\n".join(text_parts) if len(text_parts) > 1 else "[No data found in Excel file]"
    except Exception as e:
        return f"[Excel text extraction failed: {str(e)}]"


def extract_pptx(html_content: bytes) -> str:
    """Extract text from PPTX bytes using python-pptx."""
    try:
        import pptx
        import io

        prs = pptx.Presentation(io.BytesIO(html_content))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"Slide {slide_num}:"
            text_parts.append(slide_text)
            text_parts.append("=" * len(slide_text))

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

            text_parts.append("")

        return "\n\n".join(text_parts) if len(text_parts) > 1 else "[No text content found in presentation]"
    except Exception as e:
        return f"[PPTX text extraction failed: {str(e)}]"


def extract_rtf(html_content: bytes) -> str:
    """Extract text from RTF bytes (simple extraction)."""
    try:
        # Simple RTF parser - strip RTF control words
        text = html_content.decode('utf-8', errors='ignore')

        # Remove RTF control words
        rtf_pattern = r'\\[a-z]+(?:\d+)?[ ]?|\\[{}]|\\[\\\'\"]'
        plain_text = re.sub(rtf_pattern, '', text, flags=re.IGNORECASE)

        # Clean up
        plain_text = plain_text.replace('\\', ' ').strip()

        return plain_text if plain_text else "[No text content found in RTF]"
    except Exception as e:
        return f"[RTF text extraction failed: {str(e)}]"


def extract_markdown(html_content: bytes) -> str:
    """Return markdown content as-is."""
    try:
        return html_content.decode('utf-8', errors='ignore')
    except Exception:
        return "[Could not decode markdown content]"


def extract_csv(html_content: bytes) -> str:
    """Extract text from CSV bytes."""
    try:
        import csv
        import io

        text = html_content.decode('utf-8', errors='ignore')
        reader = csv.reader(io.StringIO(text))

        text_parts = []
        for row in reader:
            text_parts.append(" | ".join(row))

        return "\n".join(text_parts) if text_parts else "[No data found in CSV]"
    except Exception as e:
        return f"[CSV text extraction failed: {str(e)}]"


def extract_text(html_content: bytes) -> str:
    """Extract plain text from bytes."""
    try:
        return html_content.decode('utf-8', errors='ignore')
    except Exception:
        return "[Could not decode text content]"


def extract_document(content_bytes: bytes, content_type: str, url: str) -> Dict[str, Any]:
    """
    Extract text from document bytes based on content-type or URL.

    Args:
        content_bytes: Raw document content
        content_type: Content-Type header value
        url: Original URL (for extension detection)

    Returns:
        Dict with 'text', 'document_type', 'success' fields
    """
    # Determine document type
    doc_type = None

    # First check by URL extension
    if is_document_url(url):
        ext = Path(url.lower()).suffix
        type_map = {
            '.pdf': 'pdf',
            '.doc': 'doc', '.docx': 'docx',
            '.xls': 'xls', '.xlsx': 'xlsx',
            '.ppt': 'ppt', '.pptx': 'pptx',
            '.rtf': 'rtf',
            '.txt': 'txt',
            '.md': 'md',
            '.csv': 'csv',
            '.odt': 'odt',
            '.ods': 'ods',
            '.odp': 'odp',
        }
        doc_type = type_map.get(ext)

    # Override with content-type if available
    if content_type:
        content_type_type = get_content_type_mime(content_type)
        if content_type_type:
            doc_type = content_type_type

    # Extract based on document type
    if doc_type:
        extractors = {
            'pdf': extract_pdf,
            'doc': extract_docx,  # Fall back to docx for .doc
            'docx': extract_docx,
            'xls': extract_excel,
            'xlsx': extract_excel,
            'ppt': extract_pptx,
            'pptx': extract_pptx,
            'rtf': extract_rtf,
            'txt': extract_text,
            'md': extract_markdown,
            'csv': extract_csv,
        }

        extractor = extractors.get(doc_type, extract_text)
        text = extractor(content_bytes)

        return {
            'text': text,
            'document_type': doc_type,
            'success': '[No text content found' not in text and '[PDF text extraction failed' not in text
        }

    # If no document type detected, return as plain text
    return {
        'text': extract_text(content_bytes),
        'document_type': 'unknown',
        'success': True
    }
